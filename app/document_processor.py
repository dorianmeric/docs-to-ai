import pymupdf  # PyMuPDF -- Library for PDF processing
from pathlib import Path
from typing import List, Dict, Optional, Any
import hashlib
import json
import sys
import re
from docx import Document as DocxDocument
import pandas as pd
from bs4 import BeautifulSoup
from .config import (
    CHUNK_SIZE, CHUNK_OVERLAP, DOC_CACHE_DIR, USE_FOLDER_AS_TOPIC, DEFAULT_TOPIC, 
    CHUNKING_STRATEGY, CHUNK_BY_TOKEN, TOKENIZER_MODEL, PRESERVE_HEADINGS, 
    MAX_HEADING_CHUNK_SIZE, MIN_CHUNK_SIZE
)

try:
    import tiktoken
    TIKTOKEN_AVAILABLE = True
except ImportError:
    TIKTOKEN_AVAILABLE = False


class DocumentProcessor:
    """Handles document text extraction and chunking for PDFs and Word documents."""
    
    def __init__(self):
        self.cache_dir = DOC_CACHE_DIR
    
    def extract_topics_from_path(self, doc_path: Path, base_dir: Optional[Path] = None) -> List[str]:
        """
        Extract all topics from folder hierarchy.
        
        Args:
            doc_path: Path to the document file
            base_dir: Base directory for documents (to determine topic hierarchy)
            
        Returns:
            List of topic names from folder hierarchy
        """
        if not USE_FOLDER_AS_TOPIC:
            return [DEFAULT_TOPIC]
        
        # If no base directory, just use parent folder
        if not base_dir:
            parent_name = doc_path.parent.name
            if parent_name and parent_name not in ['/', '\\\\', '.']:
                return [parent_name]
            return [DEFAULT_TOPIC]
        
        base_path = Path(base_dir).resolve()
        doc_path_resolved = doc_path.resolve()
        
        # If document is directly in base directory
        if doc_path_resolved.parent == base_path:
            return [DEFAULT_TOPIC]
        
        # Get relative path from base to document's parent folder
        try:
            relative_path = doc_path_resolved.parent.relative_to(base_path)
        except ValueError:
            # Document is not under base directory
            return [doc_path.parent.name if doc_path.parent.name else DEFAULT_TOPIC]
        
        # Extract all folder names in the path as topics
        topics = [part for part in relative_path.parts if part and part not in ['.', '..']]
        
        return topics if topics else [DEFAULT_TOPIC]
    
    def extract_text_from_pdf(self, pdf_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text from PDF, maintaining page information.
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of dicts with 'page' and 'text'
        """
        pages_data = []
        
        try:
            doc = pymupdf.open(pdf_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                
                pages_data.append({
                    'page': page_num + 1,
                    'text': text,
                    'total_pages': len(doc)
                })
            
            doc.close()
            
        except Exception as e:
            print(f"Error processing PDF {pdf_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data
    
    def extract_text_from_docx(self, docx_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text from Word document.
        
        Args:
            docx_path: Path to the DOCX file
            
        Returns:
            List of dicts with 'page' (paragraph number) and 'text'
        """
        pages_data = []

        try:
            doc = DocxDocument(str(docx_path))
            
            # Combine paragraphs into chunks (simulating pages)
            # Group every ~10 paragraphs as a "page"
            paragraphs_per_page = 10
            current_text = []
            page_num = 1
            
            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    current_text.append(text)
                
                # Create a "page" every N paragraphs or at the end
                if (i + 1) % paragraphs_per_page == 0 or i == len(doc.paragraphs) - 1:
                    if current_text:
                        pages_data.append({
                            'page': page_num,
                            'text': '\n'.join(current_text),
                            'total_pages': -1  # Unknown for Word docs
                        })
                        page_num += 1
                        current_text = []
            
        except Exception as e:
            print(f"Error processing Word document {docx_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data
    
    def extract_text_from_markdown(self, md_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text from markdown file.
        
        Args:
            md_path: Path to the markdown file
            
        Returns:
            List of dicts with 'page' and 'text'
        """
        pages_data = []
        
        try:
            # Read the markdown file
            with open(md_path, 'r', encoding='utf-8') as f:
                markdown_content = f.read()
            
            # Convert markdown to plain text (removing markdown syntax)
            # We'll use a simple approach that preserves structure
            # For markdown files, we treat the whole file as one "page"
            # but can be split into chunks if needed
            
            # Simple conversion: remove markdown formatting
            # Remove code blocks and inline code
            
            # Remove code blocks
            text = re.sub(r'```[^`]*```', '', markdown_content, flags=re.DOTALL)
            # Remove inline code
            text = re.sub(r'`[^`]*`', '', text)
            # Remove markdown headers (## Header, ### Header, etc.)
            text = re.sub(r'^#+\s+(.*)$', r'\1', text, flags=re.MULTILINE)
            # Remove markdown lists and bullets
            text = re.sub(r'^\s*[-*+]\s+', '', text, flags=re.MULTILINE)
            # Remove markdown links
            text = re.sub(r'\[([^\\]]+)\]\([^)]+\)', r'\1', text)
            # Remove markdown images
            text = re.sub(r'!\[([^\\]]*)\]\([^)]+\)', r'\1', text)
            # Replace multiple newlines with single newline
            text = re.sub(r'\n\s*\n', '\n\n', text)
            
            pages_data.append({
                'page': 1,
                'text': text.strip(),
                'total_pages': 1
            })
            
        except Exception as e:
            print(f"Error processing markdown file {md_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data
    
    def extract_text_from_excel(self, excel_path: Path) -> List[Dict[str, Any]]:
        """
        Extract text from Excel file.
        
        Args:
            excel_path: Path to the Excel file
            
        Returns:
            List of dicts with 'page' and 'text'
        """
        pages_data = []
        
        try:
            # Read all sheets in the Excel file
            sheet_names = pd.read_excel(excel_path, sheet_name=None)
            
            # Process each sheet
            for sheet_name, df in sheet_names.items():
                # Convert dataframe to text representation
                sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                
                # Add column headers
                if not df.empty:
                    # Get column names
                    columns = list(df.columns)
                    sheet_text += "Columns: " + ", ".join(str(col) for col in columns) + "\n\n"
                    
                    # Add data rows
                    for row_num, (_, row) in enumerate(df.iterrows(), start=1):
                        row_text = "\t".join(str(val) for val in row.values)
                        sheet_text += f"Row {row_num}: {row_text}\n"
                else:
                    sheet_text += "(Empty sheet)\n"
                
                pages_data.append({
                    'page': len(pages_data) + 1,
                    'text': sheet_text.strip(),
                    'total_pages': len(sheet_names)
                })
            
        except Exception as e:
            print(f"Error processing Excel file {excel_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data

    def extract_text_from_pptx(self, pptx_path: Path) -> List[Dict[str, Any]]:
        """Extract text from PowerPoint presentation."""
        pages_data = []
        
        try:
            from pptx import Presentation
            
            prs = Presentation(str(pptx_path))
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                
                if slide_text:
                    pages_data.append({
                        'page': slide_num,
                        'text': '\n'.join(slide_text),
                        'total_pages': len(prs.slides)
                    })
                else:
                    pages_data.append({
                        'page': slide_num,
                        'text': '',
                        'total_pages': len(prs.slides)
                    })
                    
        except Exception as e:
            print(f"Error processing PowerPoint {pptx_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data

    def extract_text_from_html(self, html_path: Path) -> List[Dict[str, Any]]:
        """Extract text from HTML file."""
        pages_data = []
        
        try:
            with open(html_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text(separator='\n')
            
            lines = (line.strip() for line in text.splitlines())
            chunks = []
            current_chunk = []
            
            for line in lines:
                if line:
                    current_chunk.append(line)
                    if len('\n'.join(current_chunk)) > 2000:
                        chunks.append('\n'.join(current_chunk))
                        current_chunk = []
            
            if current_chunk:
                chunks.append('\n'.join(current_chunk))
            
            for i, chunk in enumerate(chunks, start=1):
                pages_data.append({
                    'page': i,
                    'text': chunk,
                    'total_pages': len(chunks)
                })
            
            if not pages_data:
                pages_data.append({
                    'page': 1,
                    'text': '',
                    'total_pages': 1
                })
                
        except Exception as e:
            print(f"Error processing HTML file {html_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data

    def extract_text_from_txt(self, txt_path: Path) -> List[Dict[str, Any]]:
        """Extract text from plain text file."""
        pages_data = []
        
        try:
            with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            lines = content.splitlines()
            chunks = []
            current_chunk = []
            
            for line in lines:
                current_chunk.append(line)
                if len('\n'.join(current_chunk)) > 2000:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
            
            if current_chunk:
                chunks.append('\n'.join(current_chunk))
            
            for i, chunk in enumerate(chunks, start=1):
                pages_data.append({
                    'page': i,
                    'text': chunk,
                    'total_pages': len(chunks)
                })
            
            if not pages_data:
                pages_data.append({
                    'page': 1,
                    'text': '',
                    'total_pages': 1
                })
                
        except Exception as e:
            print(f"Error processing text file {txt_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data

    def extract_text_from_csv(self, csv_path: Path) -> List[Dict[str, Any]]:
        """Extract text from CSV file with proper table handling."""
        pages_data = []
        
        try:
            df = pd.read_csv(csv_path)
            
            if not df.empty:
                csv_text = f"Columns: {', '.join(str(col) for col in df.columns)}\n\n"
                
                for idx, row in df.iterrows():
                    row_text = " | ".join(f"{col}: {row[col]}" for col in df.columns)
                    csv_text += f"Row {idx + 1}: {row_text}\n"
                
                pages_data.append({
                    'page': 1,
                    'text': csv_text.strip(),
                    'total_pages': 1
                })
            else:
                pages_data.append({
                    'page': 1,
                    'text': '(Empty CSV file)',
                    'total_pages': 1
                })
                
        except Exception as e:
            print(f"Error processing CSV file {csv_path}: {e}", file=sys.stderr)
            return []
        
        return pages_data
    
    def extract_text_from_document(self, doc_path: str, base_dir: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Extract text from document (PDF, Word, Markdown, or Excel), maintaining structure.
        
        Args:
            doc_path: Path to the document file
            base_dir: Base directory for documents (to extract topics from folder structure)
            
        Returns:
            List of dicts with 'page', 'text', and 'metadata'
        """
        doc_file = Path(doc_path)
        base_path = Path(base_dir) if base_dir else None

        # Extract topics from folder hierarchy
        topics = self.extract_topics_from_path(doc_file, base_path)

        # Get file metadata
        file_size = doc_file.stat().st_size
        last_modified = doc_file.stat().st_mtime

        # Check cache
        cache_file = self._get_cache_path(doc_file)
        if cache_file.exists():
            cached_data = json.load(open(cache_file, 'r', encoding='utf-8'))
            # Update topics in cached data in case folder structure changed
            for page_data in cached_data:
                page_data['metadata']['topics'] = topics
                page_data['metadata']['file_size'] = file_size
                page_data['metadata']['last_modified'] = last_modified
            return cached_data

        # Determine file type and extract text
        extension = doc_file.suffix.lower()

        if extension == '.pdf':
            pages_data = self.extract_text_from_pdf(doc_file)
        elif extension in ['.docx', '.doc']:
            pages_data = self.extract_text_from_docx(doc_file)
        elif extension == '.md':
            pages_data = self.extract_text_from_markdown(doc_file)
        elif extension in ['.xlsx', '.xls', '.xlsam', '.xlsb']:
            pages_data = self.extract_text_from_excel(doc_file)
        elif extension == '.pptx':
            pages_data = self.extract_text_from_pptx(doc_file)
        elif extension in ['.html', '.htm']:
            pages_data = self.extract_text_from_html(doc_file)
        elif extension == '.txt':
            pages_data = self.extract_text_from_txt(doc_file)
        elif extension == '.csv':
            pages_data = self.extract_text_from_csv(doc_file)
        else:
            print(f"Unsupported file type: {extension}", file=sys.stderr)
            return []


        print(f"\n⚠ File: {doc_file}, Found file_size = {file_size} bytes, last_modified = {last_modified} (after caching check)", file=sys.stderr)

        # Add metadata to each page
        result = []
        for page_data in pages_data:
            result.append({
                'page': page_data['page'],
                'text': page_data['text'],
                'metadata': {
                    'filename': doc_file.name,
                    'filepath': str(doc_file),
                    'topics': topics,
                    'filetype': extension,
                    'file_size': file_size,
                    'last_modified': last_modified,
                    'total_pages': page_data['total_pages']
                }
            })

        # Cache the result
        self._cache_extracted_text(doc_file, result)
        
        return result
    
    def _create_chunk(self, text: str, metadata: Dict[str, Any], page_num: int, chunk_index: int) -> Dict[str, Any]:
        """Helper to create a chunk dictionary with a unique ID."""
        topics_str = '-'.join(metadata['topics'])
        chunk_id = hashlib.md5(
            f"{topics_str}-{metadata['filepath']}-{page_num}-{chunk_index}".encode()
        ).hexdigest()
        
        return {
            'id': chunk_id,
            'text': text,
            'metadata': {
                **metadata,
                'page': page_num,
                'chunk_index': chunk_index,
            }
        }

    def _chunk_with_fixed_size(self, text: str, metadata: Dict[str, Any], page_num: int) -> List[Dict[str, Any]]:
        """Split text into fixed-size chunks with overlap."""
        chunks = []
        if not text.strip():
            return []

        for i in range(0, len(text), CHUNK_SIZE - CHUNK_OVERLAP):
            chunk_text = text[i:i + CHUNK_SIZE]
            if len(chunk_text.strip()) < 50:
                continue
            
            chunk = self._create_chunk(chunk_text, metadata, page_num, i)
            chunk['metadata']['chunk_start'] = i
            chunk['metadata']['chunk_end'] = i + len(chunk_text)
            chunks.append(chunk)
            
        return chunks

    def _chunk_by_paragraph(self, text: str, metadata: Dict[str, Any], page_num: int) -> List[Dict[str, Any]]:
        """Split text by paragraphs and group them into chunks."""
        chunks = []
        if not text.strip():
            return []

        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        
        current_chunk = ""
        chunk_index = 0
        for i, p in enumerate(paragraphs):
            # If a paragraph is larger than the chunk size, split it
            if len(p) > CHUNK_SIZE:
                if current_chunk:
                    chunks.append(self._create_chunk(current_chunk, metadata, page_num, chunk_index))
                    chunk_index += 1
                    current_chunk = ""
                
                # Split the large paragraph
                sub_chunks = self._chunk_with_fixed_size(p, metadata, page_num)
                chunks.extend(sub_chunks)
                chunk_index += len(sub_chunks)
                continue

            # Check if adding the next paragraph exceeds chunk size
            if len(current_chunk) + len(p) + 1 > CHUNK_SIZE:
                if current_chunk:
                    chunks.append(self._create_chunk(current_chunk, metadata, page_num, chunk_index))
                    chunk_index += 1
                current_chunk = p
            else:
                if current_chunk:
                    current_chunk += "\n" + p
                else:
                    current_chunk = p
        
        # Add the last remaining chunk
        if current_chunk:
            chunks.append(self._create_chunk(current_chunk, metadata, page_num, chunk_index))
            
        return chunks

    def _chunk_by_token(self, text: str, metadata: Dict[str, Any], page_num: int) -> List[Dict[str, Any]]:
        """Split text into token-based chunks with overlap."""
        chunks = []
        if not text.strip():
            return []
        
        if not TIKTOKEN_AVAILABLE:
            return self._chunk_with_fixed_size(text, metadata, page_num)
        
        try:
            enc = tiktoken.get_encoding(TOKENIZER_MODEL)
        except KeyError:
            enc = tiktoken.get_encoding("cl100k_base")
        
        tokens = enc.encode(text)
        
        for i in range(0, len(tokens), CHUNK_SIZE - CHUNK_OVERLAP):
            chunk_tokens = tokens[i:i + CHUNK_SIZE]
            chunk_text = enc.decode(chunk_tokens)
            
            if len(chunk_text.strip()) < MIN_CHUNK_SIZE:
                continue
            
            chunk = self._create_chunk(chunk_text, metadata, page_num, i)
            chunk['metadata']['chunk_start'] = i
            chunk['metadata']['chunk_end'] = i + len(chunk_tokens)
            chunk['metadata']['is_token_based'] = True
            chunks.append(chunk)
        
        return chunks

    def _extract_heading_structure(self, text: str) -> List[Dict[str, Any]]:
        """Extract heading structure from text (markdown-style or HTML)."""
        sections = []
        
        heading_pattern = re.compile(r'^(#{1,6})\s+(.+)$|^(.+?)\n(=+|-+)$', re.MULTILINE)
        
        lines = text.split('\n')
        current_section = {'heading': '', 'content': [], 'level': 0}
        
        for line in lines:
            heading_match = heading_pattern.match(line)
            
            if heading_match:
                if current_section['content']:
                    sections.append(current_section)
                    current_section = {'heading': '', 'content': [], 'level': 0}
                
                if heading_match.group(1):
                    level = len(heading_match.group(1))
                    heading = heading_match.group(2)
                else:
                    heading = heading_match.group(3)
                    level = 1 if heading_match.group(4) == '=' else 2
                
                current_section = {'heading': heading.strip(), 'content': [], 'level': level}
            else:
                current_section['content'].append(line)
        
        if current_section['content']:
            sections.append(current_section)
        
        return sections

    def _chunk_by_heading(self, text: str, metadata: Dict[str, Any], page_num: int) -> List[Dict[str, Any]]:
        """Split text by headings/sections while preserving structure."""
        chunks = []
        if not text.strip():
            return []
        
        sections = self._extract_heading_structure(text)
        
        if not sections:
            return self._chunk_with_fixed_size(text, metadata, page_num)
        
        chunk_index = 0
        for section in sections:
            section_text = '\n'.join(section['content'])
            heading = section['heading']
            
            if not section_text.strip():
                continue
            
            if len(section_text) <= MAX_HEADING_CHUNK_SIZE:
                chunk = self._create_chunk(section_text, metadata, page_num, chunk_index)
                chunk['metadata']['heading'] = heading
                chunk['metadata']['heading_level'] = section['level']
                chunk['metadata']['is_semantic_chunk'] = True
                chunks.append(chunk)
                chunk_index += 1
            else:
                sub_chunks = self._chunk_with_fixed_size(section_text, metadata, page_num)
                for sc in sub_chunks:
                    sc['metadata']['heading'] = heading
                    sc['metadata']['heading_level'] = section['level']
                    sc['metadata']['is_semantic_chunk'] = True
                chunks.extend(sub_chunks)
                chunk_index += len(sub_chunks)
        
        return chunks

    def chunk_text(self, pages_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Split text into chunks based on the configured strategy.
        
        Args:
            pages_data: List of page data from extract_text_from_document
            
        Returns:
            List of chunks with metadata
        """
        all_chunks = []
        
        use_token_chunking = CHUNK_BY_TOKEN or CHUNKING_STRATEGY == 'by_token'
        
        for page_data in pages_data:
            text = page_data['text']
            metadata = page_data['metadata']
            page_num = page_data['page']

            if CHUNKING_STRATEGY == 'by_paragraph':
                chunks = self._chunk_by_paragraph(text, metadata, page_num)
            elif CHUNKING_STRATEGY == 'semantic_heading' or (PRESERVE_HEADINGS and CHUNKING_STRATEGY != 'by_token'):
                chunks = self._chunk_by_heading(text, metadata, page_num)
            elif use_token_chunking:
                chunks = self._chunk_by_token(text, metadata, page_num)
            else:
                chunks = self._chunk_with_fixed_size(text, metadata, page_num)
            
            all_chunks.extend(chunks)
            
        return all_chunks
    
    def process_document(self, doc_path: str, base_dir: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Complete pipeline: extract text and chunk it.
        
        Args:
            doc_path: Path to document file (PDF, Word, Markdown, or Excel)
            base_dir: Base directory for documents (to extract topics from folder structure)
            
        Returns:
            List of text chunks with metadata
        """
        pages_data = self.extract_text_from_document(doc_path, base_dir)
        chunks = self.chunk_text(pages_data)
        return chunks
    
    def clear_document_cache(self):
        """Clear all cached document extractions (deletes contents, keeps folder)."""
        import shutil
        if self.cache_dir.exists():
            # Delete all files and subdirectories inside the cache folder
            for item in self.cache_dir.iterdir():
                if item.is_file():
                    item.unlink()
                elif item.is_dir():
                    shutil.rmtree(item)
            print("Document cache cleared (contents deleted, folder preserved)", file=sys.stderr)
    
    def _get_cache_path(self, doc_path: Path) -> Path:
        """Generate cache file path for a document."""
        doc_hash = hashlib.md5(str(doc_path).encode()).hexdigest()
        return self.cache_dir / f"{doc_hash}.json"
    
    def _cache_extracted_text(self, doc_path: Path, pages_data: List[Dict]):
        """Cache extracted text to avoid reprocessing."""
        cache_file = self._get_cache_path(doc_path)
        with open(cache_file, 'w', encoding='utf-8') as f:
            json.dump(pages_data, f, ensure_ascii=False, indent=2)


if __name__ == "__main__":
    # Test the processor
    processor = DocumentProcessor()
    
    # Example usage
    test_doc = "test.pdf"  # Replace with actual document path
    if Path(test_doc).exists():
        chunks = processor.process_document(test_doc)
        print(f"Extracted {len(chunks)} chunks from {test_doc}", file=sys.stderr)
        if chunks:
            print(f"\nFirst chunk preview:")
            print(f"Topics: {chunks[0]['metadata']['topics']}")
            print(chunks[0]['text'][:200])
    else:
        print(f"Test file {test_doc} not found")

