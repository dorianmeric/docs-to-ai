"""
PPTX Extractor - Extract text from PowerPoint presentations.
"""

from pathlib import Path
from typing import List, Dict, Any


def extract_text_from_pptx(pptx_path: Path) -> List[Dict[str, Any]]:
    """Extract text from PowerPoint presentation."""
    pages_data = []
    
    try:
        from pptx import Presentation
        
        prs = Presentation(str(pptx_path))
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                pages_data.append({
                    'page': slide_num,
                    'text': '\n'.join(slide_text),
                    'total_pages': len(prs.slides)
                })
            else:
                pages_data.append({
                    'page': slide_num,
                    'text': '',
                    'total_pages': len(prs.slides)
                })
                
    except Exception as e:
        print(f"Error processing PowerPoint {pptx_path}: {e}")
        return []
    
    return pages_data
